"""python-pptx renderer for SolarReach pitch decks.

Output: 14 slides, 16:9, dark theme with brand-color accents.

Slide map (must remain stable — A2 frontend embeds these in order):
  1. Cover (title + brand strip)
  2. Problem
  3. Solution (with product_description if available)
  4. Solar metrics (2x2 KPI grid: panels, kWh, saving, payback)
  5. ROI chart (25-yr cumulative cash-flow)
  6. Grid independence (large % metric)
  7. Funding comparison (5x5 table with recommended row highlighted)
  8. UK tax breaks (SEG, AIA, 0% VAT, ECO4)
  9. Decision-maker callout ("Why <role> @ <owner>")
 10. Solar radiance (flux overlay + panel grid if present)
 11. Timeline
 12. Social impact (CO₂ + trees)
 13. Tech specs
 14. CTA + footer

The renderer is defensive: any missing top-level section becomes a placeholder
slide rather than crashing the API request.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from .charts import roi_chart


# Layout constants — 16:9 widescreen
SLIDE_W_EMU = 12192000  # 13.333 in
SLIDE_H_EMU = 6858000   # 7.5 in
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Default palette (dark, deck-neutral). Brand overrides via client_doc.branding.
DEFAULT_PRIMARY = "#0F172A"   # deep navy background
DEFAULT_ACCENT = "#34D399"    # brand stripe / numerals
PANEL_BG = "#1E293B"          # card surface
PANEL_BG_ALT = "#162133"      # alt body tint
TEXT_LIGHT = "#F8FAFC"
TEXT_MUTED = "#94A3B8"
TEXT_DIM = "#64748B"
RECOMMENDED_HIGHLIGHT = "#FCD34D"  # amber for recommended row


# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(h: str) -> RGBColor:
    """Tolerant hex → RGBColor. Falls back to default primary on bad input."""
    if not isinstance(h, str):
        h = DEFAULT_PRIMARY
    h = h.lstrip("#")
    if len(h) != 6:
        h = DEFAULT_PRIMARY.lstrip("#")
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return RGBColor(0x0F, 0x17, 0x2A)


def _shade_hex(h: str, factor: float) -> str:
    """Lighten (factor>1) or darken (factor<1) a hex colour by multiplicative scale."""
    h = (h or DEFAULT_PRIMARY).lstrip("#")
    if len(h) != 6:
        h = DEFAULT_PRIMARY.lstrip("#")
    r = max(0, min(255, int(int(h[0:2], 16) * factor)))
    g = max(0, min(255, int(int(h[2:4], 16) * factor)))
    b = max(0, min(255, int(int(h[4:6], 16) * factor)))
    return f"#{r:02X}{g:02X}{b:02X}"


# ---------------------------------------------------------------------------
# Slide primitives
# ---------------------------------------------------------------------------

def _set_bg(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)


def _add_text_box(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color_hex: str = TEXT_LIGHT,
    align: str = "left",
    anchor: str = "top",
    font_name: str | None = None,
) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    run.font.color.rgb = _hex_to_rgb(color_hex)


def _add_bullets(
    slide,
    bullets: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    line_spacing: float = 1.4,
    color_hex: str = TEXT_LIGHT,
    bullet_char: str = "▸",
    bullet_color: str | None = None,
) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    bullet_color = bullet_color or color_hex
    for i, b in enumerate(bullets or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # Bullet glyph (coloured)
        r1 = p.add_run()
        r1.text = f"{bullet_char}  "
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = _hex_to_rgb(bullet_color)
        # Body text
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = _hex_to_rgb(color_hex)


def _add_rect(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_hex: str,
    line_hex: str | None = None,
    shape=MSO_SHAPE.RECTANGLE,
):
    rect = slide.shapes.add_shape(
        shape, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if line_hex:
        rect.line.color.rgb = _hex_to_rgb(line_hex)
        rect.line.width = Pt(0.75)
    else:
        rect.line.fill.background()
    return rect


def _add_accent_strip(slide, color_hex: str) -> None:
    """Brand strip down the left edge of every body slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H_EMU)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    bar.line.fill.background()


def _add_top_accent(slide, color_hex: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_EMU, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    bar.line.fill.background()


def _add_footer(slide, *, lead_id: str | None, brand_name: str, accent: str) -> None:
    _add_text_box(
        slide,
        f"Generated by SolarReach  ·  solarreach://generated/{lead_id or '—'}",
        left=0.45, top=7.05, width=8.5, height=0.3,
        font_size=8, color_hex=TEXT_DIM,
    )
    if brand_name:
        _add_text_box(
            slide,
            brand_name,
            left=9.5, top=7.05, width=3.5, height=0.3,
            font_size=8, color_hex=accent, align="right",
        )


def _add_section_header(
    slide,
    heading: str,
    *,
    accent: str,
    subhead: str | None = None,
) -> None:
    """Standard 24pt header with brand accent underline."""
    _add_text_box(
        slide, heading,
        left=0.55, top=0.45, width=12.3, height=0.7,
        font_size=24, bold=True, color_hex=TEXT_LIGHT,
    )
    # Underline accent — short bar under the header
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55), Inches(1.18), Inches(0.7), Inches(0.04),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(accent)
    bar.line.fill.background()
    if subhead:
        _add_text_box(
            slide, subhead,
            left=0.55, top=1.28, width=12.3, height=0.4,
            font_size=12, color_hex=TEXT_MUTED,
        )


# ---------------------------------------------------------------------------
# Brand resolution
# ---------------------------------------------------------------------------

def _resolve_brand(
    brand: dict[str, Any] | None,
    client_doc: dict[str, Any] | None,
) -> dict[str, str]:
    """Merge brand kwarg + client_doc.branding. brand kwarg wins on conflict."""
    branding = (client_doc or {}).get("branding") or {}
    brand = brand or {}
    primary = brand.get("primary") or branding.get("primary") or DEFAULT_PRIMARY
    accent = brand.get("accent") or branding.get("accent") or DEFAULT_ACCENT
    name = (
        brand.get("name")
        or branding.get("name")
        or (client_doc or {}).get("name")
        or ""
    )
    return {
        "primary": primary,
        "accent": accent,
        "name": name,
        "panel_bg": _shade_hex(primary, 1.25) if primary else PANEL_BG,
    }


# ---------------------------------------------------------------------------
# Funding data — used by funding comparison table
# ---------------------------------------------------------------------------

# Recommendation map: premises_type → preferred funding model name.
_RECOMMENDED_FUNDING_BY_PREMISES: dict[str, str] = {
    "office": "Capital Expense",
    "retail": "Hire Purchase",
    "warehouse": "Lease Purchase",
    "leisure": "Operational Lease",
    "education": "Free Install",
    "residential": "Free Install",
}

_FUNDING_TABLE_ROWS: list[dict[str, str]] = [
    {
        "name": "Capital Expense",
        "monthly": "—",
        "term": "0 yrs",
        "ownership": "Client (day 1)",
        "seg": "Client",
        "best_for": "Cash + tax appetite",
    },
    {
        "name": "Free Install (PPA)",
        "monthly": "Discount on unit rate",
        "term": "25 yrs",
        "ownership": "Provider",
        "seg": "Provider",
        "best_for": "Zero capex",
    },
    {
        "name": "Lease Purchase",
        "monthly": "Fixed lease + balloon",
        "term": "7 yrs",
        "ownership": "Client (end)",
        "seg": "Client",
        "best_for": "Cashflow + ownership",
    },
    {
        "name": "Operational Lease",
        "monthly": "Pure rental",
        "term": "5 yrs",
        "ownership": "Provider",
        "seg": "Client",
        "best_for": "Off-balance-sheet",
    },
    {
        "name": "Hire Purchase",
        "monthly": "Installments",
        "term": "5 yrs",
        "ownership": "Client (end)",
        "seg": "Client",
        "best_for": "Year-1 allowances",
    },
]


# UK tax / incentive lifecycle reference data (per slide 8)
_UK_TAX_BREAKS: list[dict[str, str]] = [
    {
        "name": "Smart Export Guarantee (SEG)",
        "eligibility": "Any < 5 MW system on a UK SEG-licensed network",
        "value": "≈ £0.15 / kWh exported (varies by licensee)",
        "lifecycle": "≈ £4–9k over 25 yrs at 30% export ratio",
    },
    {
        "name": "Annual Investment Allowance (AIA)",
        "eligibility": "Limited companies + sole traders, capex ≤ £1M / yr",
        "value": "100% capital allowance year 1",
        "lifecycle": "≈ 25% capex back via corp tax (£1k per £4k capex)",
    },
    {
        "name": "0% VAT (residential)",
        "eligibility": "Domestic dwellings — UK-wide until 31 Mar 2027",
        "value": "0% VAT on supply + install",
        "lifecycle": "≈ 20% off all-in cost vs. commercial",
    },
    {
        "name": "ECO4",
        "eligibility": "Low-income households / EPC D–G dwellings",
        "value": "Grant covers up to 100% of capex",
        "lifecycle": "Domestic only — confirm via supplier obligation",
    },
]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def render_pptx(
    deck_json: dict[str, Any],
    brand: dict[str, Any] | None = None,
    *,
    lead_id: str | None = None,
    out_dir: Path | str = "/tmp/decks",
    lead: dict[str, Any] | None = None,
    client_doc: dict[str, Any] | None = None,
) -> Path:
    """Render a 14-slide pitch deck to .pptx.

    Optional `lead` and `client_doc` enable real data pull-through:
      - client_doc.branding.primary / accent — slide colours
      - client_doc.product_description — Solution slide tail
      - client_doc.product_features[] — Solution slide bullets
      - client_doc.pricing_tiers[] — appendix-style rows on Solar metrics
      - lead.address, lead.owner.company_name — cover slide subline
      - lead.decision_maker.{name,role,confidence} — DM callout slide
      - lead.flux_overlay.url — embedded inferno PNG (Solar radiance slide)
      - lead.panel_layout.{panel_count, panels[]} — simple panel grid render
      - lead.premises_type — drives recommended funding row highlight
    """
    deck_json = deck_json or {}
    lead = lead or {}
    client_doc = client_doc or {}
    palette = _resolve_brand(brand, client_doc)
    primary = palette["primary"]
    accent = palette["accent"]
    brand_name = palette["name"]
    panel_bg = palette["panel_bg"]

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    fname = f"{lead_id or 'pitch'}.pptx"
    out_path = out_dir / fname

    pres = Presentation()
    pres.slide_width = SLIDE_W_EMU
    pres.slide_height = SLIDE_H_EMU
    blank = pres.slide_layouts[6]  # fully blank layout

    # ------------------------------------------------------------------ 1. Cover
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    # Big brand block on left edge (full-height accent strip)
    _add_rect(s, left=0, top=0, width=0.45, height=SLIDE_H_IN, fill_hex=accent)
    # Wordmark (top-right)
    _add_text_box(
        s,
        (brand_name or "SolarReach").upper(),
        left=8.5, top=0.5, width=4.4, height=0.4,
        font_size=12, bold=True, color_hex=accent, align="right",
    )

    title = (deck_json.get("title") or {})
    address = (lead.get("address") or "").strip()
    company = ((lead.get("owner") or {}).get("company_name") or "").strip()

    # Cover headline: prefer lead address-driven proposal line if we have it
    if company and address:
        cover_top_line = f"{company}"
        cover_bottom_line = f"Solar Proposal · {address}"
    elif address:
        cover_top_line = "Solar Proposal"
        cover_bottom_line = address
    else:
        cover_top_line = title.get("subhead", "Solar Proposal") or "Solar Proposal"
        cover_bottom_line = title.get("headline", "") or ""

    _add_text_box(
        s, cover_top_line,
        left=1.0, top=2.2, width=11.3, height=0.7,
        font_size=18, color_hex=TEXT_MUTED,
    )
    _add_text_box(
        s, title.get("headline", "Grid Independence"),
        left=1.0, top=2.85, width=11.3, height=1.5,
        font_size=46, bold=True, color_hex=TEXT_LIGHT,
    )
    _add_text_box(
        s, cover_bottom_line,
        left=1.0, top=4.45, width=11.3, height=0.5,
        font_size=18, color_hex=accent,
    )

    # Headline metric on cover (big mono number) — payback if available
    roi = deck_json.get("roi") or {}
    payback = roi.get("payback_years")
    if payback is not None:
        try:
            _add_text_box(
                s, f"{float(payback):.1f}y",
                left=1.0, top=5.4, width=3.5, height=1.3,
                font_size=72, bold=True, color_hex=accent,
                font_name="Courier New",
            )
            _add_text_box(
                s, "ESTIMATED PAYBACK",
                left=1.0, top=6.65, width=3.5, height=0.3,
                font_size=10, color_hex=TEXT_MUTED, bold=True,
            )
        except (TypeError, ValueError):
            pass

    dm = title.get("decision_maker") or ""
    if dm:
        _add_text_box(
            s, dm,
            left=8.5, top=6.45, width=4.4, height=0.5,
            font_size=14, color_hex=TEXT_LIGHT, align="right",
        )
    _add_text_box(
        s, "PREPARED BY SOLARREACH",
        left=8.5, top=6.85, width=4.4, height=0.3,
        font_size=8, color_hex=TEXT_DIM, align="right", bold=True,
    )

    # ------------------------------------------------------------------ 2. Problem
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    p = deck_json.get("problem") or {}
    _add_section_header(
        s, p.get("heading", "The grid problem"),
        accent=accent,
        subhead="UK commercial energy: volatility is the new normal",
    )
    _add_bullets(
        s, p.get("bullets", []),
        left=0.55, top=2.0, width=12.3, height=4.5,
        font_size=16, bullet_color=accent,
    )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 3. Solution
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    sol = deck_json.get("solution") or {}
    _add_section_header(
        s, sol.get("heading", "Own your generation"),
        accent=accent,
        subhead="Roof-mounted PV: 30-year asset, fixed marginal cost",
    )

    # Use Sonnet's bullets if present, plus product_features from client_doc
    bullets = list(sol.get("bullets") or [])
    pfeats = client_doc.get("product_features") or []
    if isinstance(pfeats, list):
        for f in pfeats[:4]:
            if isinstance(f, str):
                if f not in bullets:
                    bullets.append(f)
            elif isinstance(f, dict):
                label = f.get("name") or f.get("label") or ""
                if label and label not in bullets:
                    bullets.append(label)
    _add_bullets(
        s, bullets[:6],
        left=0.55, top=2.0, width=7.5, height=3.8,
        font_size=14, bullet_color=accent,
    )

    # Right-side: product_description card if available
    pdesc = (client_doc.get("product_description") or "").strip()
    if pdesc:
        _add_rect(
            s, left=8.4, top=2.0, width=4.5, height=4.0,
            fill_hex=panel_bg, line_hex=accent,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        _add_text_box(
            s, "WHAT WE INSTALL",
            left=8.65, top=2.15, width=4.0, height=0.3,
            font_size=9, bold=True, color_hex=accent,
        )
        _add_text_box(
            s, pdesc[:480],
            left=8.65, top=2.55, width=4.0, height=3.3,
            font_size=11, color_hex=TEXT_LIGHT,
        )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 4. Solar metrics (KPI grid)
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    _add_section_header(
        s, "Solar metrics for this site",
        accent=accent,
        subhead=address or "From Google Solar API + UK irradiance modelling",
    )

    pl = (lead.get("panel_layout") or {})
    fin = (lead.get("financial") or {})
    tech = deck_json.get("tech_specs") or {}

    panel_count = pl.get("panel_count") or tech.get("panels") or 0
    annual_kwh = pl.get("annual_kwh") or tech.get("annual_kwh") or roi.get("annual_kwh") or 0
    annual_saving = roi.get("annual_saving_gbp") or fin.get("annual_saving_gbp") or 0
    payback_y = roi.get("payback_years") or fin.get("payback_years") or 0
    npv = roi.get("npv_25yr_gbp") or fin.get("npv_25yr_gbp") or 0

    # 2x2 KPI grid
    kpis = [
        ("PANELS",            f"{int(panel_count or 0)}",       "On usable rooftop"),
        ("ANNUAL GENERATION", f"{int(annual_kwh or 0):,} kWh",  "Year 1, pre-degradation"),
        ("ANNUAL SAVING",     f"£{int(annual_saving or 0):,}",  "vs. grid at current tariff"),
        ("PAYBACK",           f"{float(payback_y or 0):.1f}y",  "Cumulative cash positive"),
    ]
    grid_left = 0.55
    grid_top = 2.0
    cell_w = 6.15
    cell_h = 2.0
    for idx, (label, value, sub) in enumerate(kpis):
        col = idx % 2
        row = idx // 2
        x = grid_left + col * (cell_w + 0.1)
        y = grid_top + row * (cell_h + 0.15)
        # Card background
        _add_rect(
            s, left=x, top=y, width=cell_w, height=cell_h,
            fill_hex=panel_bg,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        # Brand-accented underline at the bottom of each card
        _add_rect(
            s, left=x + 0.2, top=y + cell_h - 0.18,
            width=0.7, height=0.05, fill_hex=accent,
        )
        # Label
        _add_text_box(
            s, label,
            left=x + 0.2, top=y + 0.15, width=cell_w - 0.4, height=0.35,
            font_size=10, bold=True, color_hex=TEXT_MUTED,
        )
        # Big mono number
        _add_text_box(
            s, value,
            left=x + 0.2, top=y + 0.45, width=cell_w - 0.4, height=1.0,
            font_size=40, bold=True, color_hex=TEXT_LIGHT,
            font_name="Courier New",
        )
        # Sub-label
        _add_text_box(
            s, sub,
            left=x + 0.2, top=y + 1.5, width=cell_w - 0.4, height=0.35,
            font_size=10, color_hex=accent,
        )

    # NPV banner across bottom
    _add_text_box(
        s,
        f"25-YEAR NPV  ·  £{int(npv or 0):,}",
        left=0.55, top=6.4, width=12.3, height=0.4,
        font_size=14, bold=True, color_hex=accent, align="center",
    )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 5. ROI chart
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    _add_section_header(
        s, roi.get("heading", "Return on investment"),
        accent=accent,
        subhead="25-year cumulative cash flow with 0.5%/yr panel degradation",
    )

    capex = float(roi.get("capex_gbp") or fin.get("capex_gbp") or 0)
    saving_v = float(annual_saving or 0)
    payback_v = float(payback_y or 0)
    npv_v = float(npv or 0)

    # Render + embed chart, full-bleed below the header
    try:
        chart_path = roi_chart(
            payback_years=payback_v,
            capex=capex,
            annual_saving=saving_v,
            npv_25yr=npv_v,
            out_dir=out_dir,
            lead_id=lead_id or "deck",
        )
        s.shapes.add_picture(
            str(chart_path),
            Inches(0.55), Inches(1.85), Inches(12.3), Inches(4.7),
        )
    except Exception:
        _add_text_box(
            s, "[ROI chart unavailable]",
            left=0.55, top=1.85, width=12.3, height=4.7,
            font_size=20, color_hex=TEXT_MUTED, align="center",
        )

    # KPI strip below the chart
    _add_text_box(
        s,
        f"Capex £{int(capex):,}    ·    Annual saving £{int(saving_v):,}    "
        f"·    Payback {payback_v:.1f}y    ·    25-yr NPV £{int(npv_v):,}",
        left=0.55, top=6.65, width=12.3, height=0.35,
        font_size=11, color_hex=TEXT_MUTED, align="center",
    )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 6. Grid independence
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    gi = deck_json.get("grid_independence") or {}
    _add_section_header(
        s, gi.get("heading", "Grid independence"),
        accent=accent,
        subhead="Your roof. Your generation. Your terms.",
    )
    _add_text_box(
        s, gi.get("body", ""),
        left=0.55, top=2.0, width=7.5, height=3.0,
        font_size=18, color_hex=TEXT_LIGHT,
    )
    pct = gi.get("metric_pct_offset")
    if pct is not None:
        try:
            _add_text_box(
                s, f"{int(pct)}%",
                left=8.5, top=1.9, width=4.4, height=2.6,
                font_size=120, bold=True, color_hex=accent,
                font_name="Courier New", align="center",
            )
            _add_text_box(
                s, "OF GRID DEMAND OFFSET ON-SITE",
                left=8.5, top=4.55, width=4.4, height=0.4,
                font_size=11, color_hex=TEXT_MUTED, align="center", bold=True,
            )
        except (TypeError, ValueError):
            pass
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 7. Funding comparison table
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    fnd = deck_json.get("funding") or {}
    _add_section_header(
        s, fnd.get("heading", "5 ways to fund it"),
        accent=accent,
        subhead="Recommended row highlighted for your premises type",
    )

    # Recommended row by premises_type
    premises = (lead.get("premises_type") or "").lower()
    recommended = _RECOMMENDED_FUNDING_BY_PREMISES.get(premises)

    # Render a manual table (rectangles) — gives us full visual control
    headers = ["Model", "Monthly cost", "Term", "Ownership at end", "SEG claimer", "Best for"]
    col_widths = [2.8, 2.0, 1.0, 2.3, 1.5, 2.7]  # in (sum ≤ 12.3)
    table_left = 0.55
    table_top = 2.0
    header_h = 0.45
    row_h = 0.7

    # Header row
    x = table_left
    for i, hcell in enumerate(headers):
        _add_rect(
            s, left=x, top=table_top, width=col_widths[i], height=header_h,
            fill_hex=accent,
        )
        _add_text_box(
            s, hcell,
            left=x + 0.1, top=table_top + 0.05, width=col_widths[i] - 0.2, height=header_h - 0.1,
            font_size=10, bold=True, color_hex=primary, align="left", anchor="middle",
        )
        x += col_widths[i]

    # Body rows
    for r_idx, row in enumerate(_FUNDING_TABLE_ROWS):
        y = table_top + header_h + r_idx * row_h
        is_recommended = recommended and row["name"] == recommended
        row_bg = RECOMMENDED_HIGHLIGHT if is_recommended else (
            panel_bg if r_idx % 2 == 0 else PANEL_BG_ALT
        )
        text_color = primary if is_recommended else TEXT_LIGHT
        muted_color = primary if is_recommended else TEXT_MUTED
        cells = [
            row["name"],
            row["monthly"],
            row["term"],
            row["ownership"],
            row["seg"],
            row["best_for"],
        ]
        x = table_left
        for c_idx, cell_text in enumerate(cells):
            _add_rect(
                s, left=x, top=y, width=col_widths[c_idx], height=row_h,
                fill_hex=row_bg,
            )
            _add_text_box(
                s, cell_text,
                left=x + 0.1, top=y + 0.1, width=col_widths[c_idx] - 0.2, height=row_h - 0.2,
                font_size=10,
                bold=(c_idx == 0),
                color_hex=text_color if c_idx == 0 else muted_color,
                align="left", anchor="middle",
            )
            x += col_widths[c_idx]
        # "RECOMMENDED" pill on highlighted row
        if is_recommended:
            _add_text_box(
                s, "★ RECOMMENDED",
                left=table_left + col_widths[0] - 1.4, top=y + 0.15,
                width=1.3, height=0.3,
                font_size=8, bold=True, color_hex=primary, align="right",
            )

    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 8. UK tax breaks
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    _add_section_header(
        s, "UK incentives & tax breaks",
        accent=accent,
        subhead="What HMRC and the regulator give back",
    )
    # 4 tax cards in a 2x2 grid
    tb_left = 0.55
    tb_top = 2.0
    card_w = 6.15
    card_h = 2.1
    for i, tb in enumerate(_UK_TAX_BREAKS):
        col = i % 2
        row = i // 2
        x = tb_left + col * (card_w + 0.1)
        y = tb_top + row * (card_h + 0.15)
        _add_rect(
            s, left=x, top=y, width=card_w, height=card_h,
            fill_hex=panel_bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        _add_rect(
            s, left=x, top=y, width=0.1, height=card_h, fill_hex=accent,
        )
        _add_text_box(
            s, tb["name"],
            left=x + 0.25, top=y + 0.15, width=card_w - 0.4, height=0.4,
            font_size=14, bold=True, color_hex=accent,
        )
        _add_text_box(
            s, f"Eligibility:  {tb['eligibility']}",
            left=x + 0.25, top=y + 0.6, width=card_w - 0.4, height=0.4,
            font_size=10, color_hex=TEXT_MUTED,
        )
        _add_text_box(
            s, f"Value:  {tb['value']}",
            left=x + 0.25, top=y + 1.0, width=card_w - 0.4, height=0.4,
            font_size=10, color_hex=TEXT_LIGHT,
        )
        _add_text_box(
            s, f"Lifecycle:  {tb['lifecycle']}",
            left=x + 0.25, top=y + 1.4, width=card_w - 0.4, height=0.5,
            font_size=10, color_hex=TEXT_LIGHT,
        )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 9. Decision-maker callout
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    dmc = deck_json.get("decision_maker_callout") or {}
    dm_lead = (lead.get("decision_maker") or {})
    dm_role = dm_lead.get("role") or "Director"
    dm_name = dm_lead.get("name") or ""
    dm_conf = dm_lead.get("confidence")
    owner_name = ((lead.get("owner") or {}).get("company_name") or "").strip()
    if owner_name:
        why_heading = f"Why {dm_role} @ {owner_name}"
    else:
        why_heading = dmc.get("heading", f"For {dm_role}")
    _add_section_header(
        s, why_heading,
        accent=accent,
        subhead=f"Tuned to the concerns of a {dm_role.lower()}",
    )
    if dm_name:
        _add_text_box(
            s, dm_name.upper(),
            left=0.55, top=1.95, width=12.3, height=0.4,
            font_size=11, bold=True, color_hex=accent,
        )
    _add_text_box(
        s, dmc.get("body", ""),
        left=0.55, top=2.4, width=8.5, height=4.0,
        font_size=16, color_hex=TEXT_LIGHT,
    )

    # Confidence card on the right
    if dm_conf is not None:
        try:
            conf_pct = int(round(float(dm_conf) * 100))
        except (TypeError, ValueError):
            conf_pct = None
        _add_rect(
            s, left=9.5, top=2.4, width=3.4, height=2.5,
            fill_hex=panel_bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        _add_text_box(
            s, "DECISION-MAKER CONFIDENCE",
            left=9.65, top=2.55, width=3.1, height=0.4,
            font_size=9, bold=True, color_hex=TEXT_MUTED,
        )
        if conf_pct is not None:
            _add_text_box(
                s, f"{conf_pct}%",
                left=9.65, top=2.95, width=3.1, height=1.4,
                font_size=64, bold=True, color_hex=accent,
                font_name="Courier New", align="center",
            )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 10. Solar radiance / panel layout
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    _add_section_header(
        s, "Solar radiance & panel layout",
        accent=accent,
        subhead="Google Solar API flux + clipped rooftop layout",
    )

    flux_url = ((lead.get("flux_overlay") or {}).get("url") or "").strip()
    embedded_flux = False
    if flux_url:
        # Try local path first (e.g. /static/flux/<id>.png mapped to /tmp/flux)
        candidate_paths = []
        if flux_url.startswith("/"):
            candidate_paths.extend([
                Path("/tmp") / flux_url.lstrip("/"),
                Path("/tmp/flux") / Path(flux_url).name,
                Path("/tmp/decks") / Path(flux_url).name,
            ])
        for cp in candidate_paths:
            try:
                if cp.exists() and cp.stat().st_size > 100:
                    s.shapes.add_picture(
                        str(cp), Inches(0.55), Inches(1.95),
                        Inches(6.0), Inches(4.6),
                    )
                    embedded_flux = True
                    break
            except Exception:
                continue
    if not embedded_flux:
        # Placeholder card if we don't have the file locally
        _add_rect(
            s, left=0.55, top=1.95, width=6.0, height=4.6,
            fill_hex=panel_bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        _add_text_box(
            s, "FLUX OVERLAY",
            left=0.7, top=2.1, width=5.7, height=0.4,
            font_size=10, bold=True, color_hex=accent,
        )
        if flux_url:
            _add_text_box(
                s, f"Available at\n{flux_url}",
                left=0.7, top=2.6, width=5.7, height=0.8,
                font_size=11, color_hex=TEXT_MUTED,
            )
        else:
            _add_text_box(
                s, "Not generated for this lead.\nRun /lead/<id>/flux_overlay to populate.",
                left=0.7, top=2.6, width=5.7, height=1.0,
                font_size=11, color_hex=TEXT_MUTED,
            )

    # Right side: simple panel grid render
    _add_rect(
        s, left=6.7, top=1.95, width=6.15, height=4.6,
        fill_hex=panel_bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    _add_text_box(
        s, "PANEL LAYOUT",
        left=6.85, top=2.1, width=5.85, height=0.4,
        font_size=10, bold=True, color_hex=accent,
    )

    panel_count_int = int(panel_count or 0)
    if panel_count_int > 0:
        # Lay out as ~4:3 grid
        import math
        cols = max(1, int(math.ceil(math.sqrt(panel_count_int * 1.4))))
        rows = max(1, int(math.ceil(panel_count_int / cols)))
        avail_w = 5.7
        avail_h = 3.4
        cell_w_panel = min(avail_w / cols, 0.42)
        cell_h_panel = min(avail_h / rows, 0.30)
        gap = 0.04
        total_w = cols * (cell_w_panel + gap)
        total_h = rows * (cell_h_panel + gap)
        ox = 6.85 + (5.85 - total_w) / 2
        oy = 2.6 + (3.6 - total_h) / 2
        drawn = 0
        for r in range(rows):
            for c in range(cols):
                if drawn >= panel_count_int:
                    break
                px = ox + c * (cell_w_panel + gap)
                py = oy + r * (cell_h_panel + gap)
                _add_rect(
                    s, left=px, top=py, width=cell_w_panel, height=cell_h_panel,
                    fill_hex=accent, shape=MSO_SHAPE.RECTANGLE,
                )
                drawn += 1
        _add_text_box(
            s,
            f"{panel_count_int} panels  ·  {int(annual_kwh or 0):,} kWh / yr",
            left=6.85, top=6.15, width=5.85, height=0.35,
            font_size=11, color_hex=TEXT_LIGHT, bold=True, align="center",
        )
    else:
        _add_text_box(
            s, "No panel layout yet — run /lead/<id>/panels.",
            left=6.85, top=2.6, width=5.85, height=1.0,
            font_size=11, color_hex=TEXT_MUTED,
        )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 11. Timeline
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    tl = deck_json.get("timeline") or {}
    _add_section_header(
        s, tl.get("heading", "Timeline to commissioning"),
        accent=accent,
        subhead="Survey → Design + DNO → Install → Commissioning",
    )
    phases = tl.get("phases") or []
    if phases:
        n = max(1, len(phases))
        col_w = (12.0) / n
        # Horizontal connector line
        line_y = 3.6
        line = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), Inches(line_y),
            Inches(12.0 - 0.2), Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_to_rgb(accent)
        line.line.fill.background()

        for i, ph in enumerate(phases):
            x = 0.7 + i * col_w
            cx = x + col_w / 2
            # Numbered circle
            dot = s.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - 0.25), Inches(line_y - 0.25 + 0.02),
                Inches(0.5), Inches(0.5),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = _hex_to_rgb(accent)
            dot.line.color.rgb = _hex_to_rgb(primary)
            dot.line.width = Pt(2)
            _add_text_box(
                s, str(i + 1),
                left=cx - 0.25, top=line_y - 0.18,
                width=0.5, height=0.4,
                font_size=14, bold=True, color_hex=primary, align="center",
            )
            _add_text_box(
                s, ph.get("name", ""),
                left=x, top=2.4, width=col_w, height=0.5,
                font_size=14, bold=True, color_hex=TEXT_LIGHT, align="center",
            )
            wks = ph.get("weeks", "?")
            _add_text_box(
                s, f"{wks} wks",
                left=x, top=4.2, width=col_w, height=0.4,
                font_size=12, color_hex=accent, align="center", bold=True,
            )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 12. Social impact
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    si = deck_json.get("social_impact") or {}
    _add_section_header(
        s, si.get("heading", "Carbon offset"),
        accent=accent,
        subhead="UK grid intensity 0.233 kg CO₂/kWh; 22 kg CO₂ / tree / year",
    )
    co2 = si.get("tonnes_co2_yr")
    trees = si.get("equiv_trees")
    if co2 is not None:
        _add_text_box(
            s, f"{co2}",
            left=0.55, top=2.2, width=6.0, height=2.4,
            font_size=120, bold=True, color_hex=accent,
            font_name="Courier New",
        )
        _add_text_box(
            s, "TONNES CO₂ AVOIDED PER YEAR",
            left=0.55, top=4.7, width=6.0, height=0.4,
            font_size=11, bold=True, color_hex=TEXT_MUTED,
        )
    if trees is not None:
        _add_rect(
            s, left=7.0, top=2.2, width=5.85, height=3.0,
            fill_hex=panel_bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        _add_text_box(
            s, "TREE EQUIVALENT",
            left=7.2, top=2.4, width=5.5, height=0.4,
            font_size=10, bold=True, color_hex=TEXT_MUTED,
        )
        _add_text_box(
            s, f"≈ {trees}",
            left=7.2, top=2.85, width=5.5, height=1.4,
            font_size=72, bold=True, color_hex=accent,
            font_name="Courier New", align="center",
        )
        _add_text_box(
            s, "trees planted, every year",
            left=7.2, top=4.4, width=5.5, height=0.4,
            font_size=12, color_hex=TEXT_LIGHT, align="center",
        )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 13. Tech specs
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    _add_accent_strip(s, accent)
    _add_top_accent(s, accent)
    ts = deck_json.get("tech_specs") or {}
    _add_section_header(
        s, ts.get("heading", "Tech specifications"),
        accent=accent,
        subhead="Standard build — Tier 1 panels, 25-year linear performance warranty",
    )
    rows = [
        ("Panels",             ts.get("panels", panel_count or "—")),
        ("Peak power (kWp)",   ts.get("kw_peak", "—")),
        ("Annual generation",  f"{int(ts.get('annual_kwh') or annual_kwh or 0):,} kWh"),
        ("Warranty (years)",   ts.get("warranty_years", 25)),
        ("Inverter",           ts.get("inverter", "Hybrid string + battery-ready")),
        ("Monitoring",         ts.get("monitoring", "API + dashboard, 1-min granularity")),
    ]
    for i, (k, v) in enumerate(rows):
        y = 2.0 + i * 0.7
        # Alternating row tint
        if i % 2 == 0:
            _add_rect(
                s, left=0.55, top=y, width=12.3, height=0.6,
                fill_hex=panel_bg,
            )
        _add_text_box(
            s, k,
            left=0.8, top=y + 0.1, width=5.0, height=0.4,
            font_size=14, color_hex=TEXT_MUTED, anchor="middle",
        )
        _add_text_box(
            s, str(v),
            left=6.0, top=y + 0.1, width=6.5, height=0.4,
            font_size=14, bold=True, color_hex=accent, anchor="middle",
        )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    # ------------------------------------------------------------------ 14. CTA
    s = pres.slides.add_slide(blank)
    _set_bg(s, primary)
    # Full-width brand block down the left third
    _add_rect(s, left=0, top=0, width=0.45, height=SLIDE_H_IN, fill_hex=accent)
    cta = deck_json.get("cta") or {}
    _add_text_box(
        s, "NEXT STEP",
        left=1.0, top=2.2, width=11.5, height=0.4,
        font_size=12, bold=True, color_hex=accent,
    )
    _add_text_box(
        s, cta.get("heading", "Book a 30-minute walkthrough"),
        left=1.0, top=2.7, width=11.5, height=1.4,
        font_size=44, bold=True, color_hex=TEXT_LIGHT,
    )
    _add_text_box(
        s, cta.get("body", ""),
        left=1.0, top=4.3, width=11.5, height=1.4,
        font_size=18, color_hex=TEXT_MUTED,
    )
    _add_text_box(
        s, cta.get("contact", ""),
        left=1.0, top=5.7, width=11.5, height=0.5,
        font_size=18, bold=True, color_hex=accent,
    )
    _add_footer(s, lead_id=lead_id, brand_name=brand_name, accent=accent)

    pres.save(str(out_path))
    return out_path
