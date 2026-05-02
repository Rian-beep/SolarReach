"""PPTX renderer tests — produces a valid 14-slide 16:9 deck.

The renderer is the user-facing surface of every pitch — these tests guard:
  - slide count + 16:9 dimensions (frontend embed contract)
  - filename = pitch_<lead_id>.pptx
  - graceful degradation on partial deck JSON
  - brand colour pull-through from client_doc.branding
  - KPI grid presence on the metrics slide
  - funding comparison table presence + recommended row highlighting
"""

from __future__ import annotations


SAMPLE_DECK = {
    "title": {
        "headline": "Grid Independence — 1 Old St",
        "subhead": "Solar PV proposal for Old Street Holdings Ltd",
        "decision_maker": "Sarah Patel, CFO",
    },
    "problem": {
        "heading": "The grid problem",
        "bullets": [
            "UK commercial electricity has risen 78% since 2019",
            "Volatility is the new normal",
        ],
    },
    "solution": {
        "heading": "Rooftop solar — own your generation",
        "bullets": ["10,080 kWh/year on-site", "24 panels on existing roof", "30-year asset"],
    },
    "grid_independence": {
        "heading": "Grid independence",
        "body": "Your roof, your generation, your terms.",
        "metric_pct_offset": 42,
    },
    "roi": {
        "heading": "ROI",
        "capex_gbp": 24500,
        "annual_saving_gbp": 3120,
        "payback_years": 7.8,
        "npv_25yr_gbp": 41200,
    },
    "funding": {
        "heading": "5 ways to fund it",
        "models": [
            {"name": "Capital Expense", "fit": "Best NPV"},
            {"name": "Free Install", "fit": "Zero capex"},
            {"name": "Lease Purchase", "fit": "Balance"},
            {"name": "Operational Lease", "fit": "Off-balance-sheet"},
            {"name": "Hire Purchase", "fit": "VAT-efficient"},
        ],
    },
    "timeline": {
        "heading": "Timeline",
        "phases": [
            {"name": "Survey", "weeks": 2},
            {"name": "Design + DNO", "weeks": 6},
            {"name": "Install", "weeks": 4},
            {"name": "Commissioning", "weeks": 1},
        ],
    },
    "decision_maker_callout": {
        "heading": "For Sarah Patel, CFO",
        "body": "This is a finance decision dressed as a roof decision.",
    },
    "social_impact": {
        "heading": "Carbon offset",
        "tonnes_co2_yr": 2.4,
        "equiv_trees": 110,
    },
    "tech_specs": {
        "heading": "Tech",
        "panels": 24,
        "kw_peak": 9.6,
        "annual_kwh": 10080,
        "warranty_years": 25,
    },
    "cta": {
        "heading": "Next step",
        "body": "30-min call this week with Sarah Patel.",
        "contact": "hello@greensolar.uk",
    },
}


SAMPLE_BRAND = {"primary": "#0F172A", "logo_url": None, "name": "GreenSolar UK"}

SAMPLE_LEAD = {
    "_id": "lead_test_xyz",
    "address": "1 Old Street, London EC1Y 8AF",
    "premises_type": "office",
    "owner": {"company_name": "Old Street Holdings Ltd"},
    "decision_maker": {"name": "Sarah Patel", "role": "CFO", "confidence": 0.92},
    "panel_layout": {"panel_count": 24, "annual_kwh": 10080},
    "financial": {
        "capex_gbp": 24500,
        "annual_saving_gbp": 3120,
        "payback_years": 7.8,
        "npv_25yr_gbp": 41200,
    },
}


SAMPLE_CLIENT_DOC = {
    "_id": "client-greensolar-uk",
    "name": "GreenSolar UK",
    "branding": {"primary": "#123456", "accent": "#ABCDEF"},
    "product_description": "Tier-1 monocrystalline panels with hybrid inverter and battery readiness.",
    "product_features": [
        "25-year linear performance warranty",
        "Sub-1-minute monitoring API",
        "Battery-ready inverter",
    ],
    "pricing_tiers": [
        {"name": "Standard", "panel_unit_gbp": 850},
    ],
}


# Total slide count after redesign (cover + 12 body + cta).
EXPECTED_SLIDES = 14


def test_render_pptx_creates_full_deck(tmp_path):
    from codex_brain.generators.pptx_renderer import render_pptx
    pptx_path = render_pptx(SAMPLE_DECK, SAMPLE_BRAND, lead_id="lead_test_xyz", out_dir=tmp_path)
    assert pptx_path.exists()
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    assert len(pres.slides) == EXPECTED_SLIDES


def test_render_pptx_is_widescreen_169(tmp_path):
    from codex_brain.generators.pptx_renderer import render_pptx
    pptx_path = render_pptx(SAMPLE_DECK, SAMPLE_BRAND, lead_id="lead_test_169", out_dir=tmp_path)
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    # 16:9 dims in EMU: 12192000 x 6858000
    assert pres.slide_width == 12192000
    assert pres.slide_height == 6858000


def test_render_pptx_filename_uses_lead_id(tmp_path):
    from codex_brain.generators.pptx_renderer import render_pptx
    pptx_path = render_pptx(SAMPLE_DECK, SAMPLE_BRAND, lead_id="lead_my_id", out_dir=tmp_path)
    assert "lead_my_id" in pptx_path.name


def test_render_pptx_handles_missing_optional_fields(tmp_path):
    """Sonnet might omit a section — render should still produce all slides with placeholders."""
    from codex_brain.generators.pptx_renderer import render_pptx
    minimal = {"title": {"headline": "Test"}}
    pptx_path = render_pptx(minimal, SAMPLE_BRAND, lead_id="lead_min", out_dir=tmp_path)
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    assert len(pres.slides) == EXPECTED_SLIDES


# ---------------------------------------------------------------------------
# New tests covering brand pull-through, KPI grid, funding table
# ---------------------------------------------------------------------------

def _all_text(pres) -> str:
    """Concatenate every text run across every slide for substring assertions."""
    out = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    out.append(run.text or "")
    return "\n".join(out)


def _slide_fill_colors(slide) -> set[str]:
    """All solid fill RGB hex values found on a slide (shapes + background)."""
    colors: set[str] = set()
    try:
        bg = slide.background.fill
        if bg.type is not None:  # solid
            rgb = bg.fore_color.rgb
            if rgb is not None:
                colors.add(str(rgb).upper())
    except Exception:
        pass
    for shape in slide.shapes:
        try:
            f = shape.fill
            if f.type is not None and f.fore_color.rgb is not None:
                colors.add(str(f.fore_color.rgb).upper())
        except Exception:
            continue
    return colors


def test_brand_colours_pull_through_from_client_doc(tmp_path):
    """client_doc.branding.primary + accent should propagate to slide fills."""
    from codex_brain.generators.pptx_renderer import render_pptx
    # Pass brand=None so client_doc.branding is the only source.
    pptx_path = render_pptx(
        SAMPLE_DECK,
        brand=None,
        lead_id="lead_brand",
        out_dir=tmp_path,
        lead=SAMPLE_LEAD,
        client_doc=SAMPLE_CLIENT_DOC,
    )
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    primary_hex = SAMPLE_CLIENT_DOC["branding"]["primary"].lstrip("#").upper()
    accent_hex = SAMPLE_CLIENT_DOC["branding"]["accent"].lstrip("#").upper()
    all_colours: set[str] = set()
    for slide in pres.slides:
        all_colours |= _slide_fill_colors(slide)
    assert primary_hex in all_colours, (
        f"client_doc.branding.primary ({primary_hex}) not used as a fill colour"
    )
    assert accent_hex in all_colours, (
        f"client_doc.branding.accent ({accent_hex}) not used as a fill colour"
    )

    # product_description should land in the deck (Solution slide right card).
    text_blob = _all_text(pres)
    assert "monocrystalline" in text_blob.lower(), (
        "product_description content not found anywhere in deck"
    )


def test_solar_metrics_kpi_grid_present(tmp_path):
    """Slide 4 (Solar metrics) must show the four KPI labels and rendered values."""
    from codex_brain.generators.pptx_renderer import render_pptx
    pptx_path = render_pptx(
        SAMPLE_DECK,
        SAMPLE_BRAND,
        lead_id="lead_kpi",
        out_dir=tmp_path,
        lead=SAMPLE_LEAD,
        client_doc=SAMPLE_CLIENT_DOC,
    )
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    # Slide index 3 is the Solar metrics slide (0-based: cover/problem/solution/metrics).
    metrics_slide = pres.slides[3]
    text = "\n".join(
        run.text
        for shape in metrics_slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    for required in ["PANELS", "ANNUAL GENERATION", "ANNUAL SAVING", "PAYBACK"]:
        assert required in text, f"KPI label '{required}' missing from metrics slide"
    # Values pulled from lead.panel_layout / lead.financial.
    assert "24" in text, "panel_count value (24) missing"
    assert "10,080" in text, "annual_kwh value (10,080) missing"
    assert "3,120" in text, "annual_saving value (3,120) missing"
    assert "7.8" in text, "payback (7.8) missing"
    assert "NPV" in text.upper(), "NPV banner missing from metrics slide"


def test_funding_comparison_table_present_and_recommended_highlighted(tmp_path):
    """Slide 7 (funding) must list all 5 models, header columns, and highlight by premises_type."""
    from codex_brain.generators.pptx_renderer import render_pptx
    pptx_path = render_pptx(
        SAMPLE_DECK,
        SAMPLE_BRAND,
        lead_id="lead_funding",
        out_dir=tmp_path,
        lead=SAMPLE_LEAD,  # premises_type=office → recommended = Capital Expense
        client_doc=SAMPLE_CLIENT_DOC,
    )
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    funding_slide = pres.slides[6]  # 0-based: index 6 is the 7th slide
    text = "\n".join(
        run.text
        for shape in funding_slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    # All five funding model names must appear in the table.
    for model in [
        "Capital Expense",
        "Free Install",
        "Lease Purchase",
        "Operational Lease",
        "Hire Purchase",
    ]:
        assert model in text, f"funding model '{model}' missing from comparison table"
    # All five column headers must appear.
    for header in ["Monthly cost", "Term", "Ownership", "SEG", "Best for"]:
        assert header in text, f"funding column '{header}' missing"
    # Recommended pill should appear for the office → Capital Expense match.
    assert "RECOMMENDED" in text, "premises_type='office' should highlight a recommended row"

    # Highlight colour (#FCD34D amber) must show up among shape fills on the slide.
    fills = _slide_fill_colors(funding_slide)
    assert "FCD34D" in fills, "Recommended-row highlight colour missing"
