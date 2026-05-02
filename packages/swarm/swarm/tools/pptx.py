"""Pitch-deck builder tool.

Tries to reuse `codex_brain.generators.pptx_renderer.render_pptx` (the existing
11-slide branded renderer). If that import fails, falls back to a minimal
inline python-pptx renderer that produces a 3-slide deck (title / body / CTA).
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from langchain_core.tools import tool

from swarm.audit import get_actor_name, write_audit_sync
from swarm.mongo import get_sync_db

log = logging.getLogger("solarreach.swarm.tools.pptx")

OUT_DIR = Path("/tmp/swarm-decks")
OUT_DIR.mkdir(parents=True, exist_ok=True)

try:
    from codex_brain.generators.pptx_renderer import render_pptx as _render_pptx_full  # type: ignore
except Exception:  # noqa: BLE001
    _render_pptx_full = None


def _render_minimal(deck_json: dict[str, Any], lead_id: str | None) -> Path:
    """Inline 3-slide fallback when codex_brain isn't importable."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pres = Presentation()
    pres.slide_width = 12192000
    pres.slide_height = 6858000
    blank = pres.slide_layouts[6]

    title = (deck_json.get("title") or {}).get("headline") or "Solar Proposal"
    subhead = (deck_json.get("title") or {}).get("subhead") or ""
    body = (deck_json.get("solution") or {}).get("bullets") or ["Your rooftop is generation-ready."]
    cta = (deck_json.get("cta") or {}).get("body") or "Let's schedule a 15-minute call."

    for txt, size, top in [(title, 44, 2.4), (subhead, 20, 4.2)]:
        s = pres.slides.add_slide(blank) if top == 2.4 else s
        tb = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(1.4))
        tb.text_frame.text = txt
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(size) if tb.text_frame.paragraphs[0].runs else Pt(size)

    s = pres.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12), Inches(0.8))
    tb.text_frame.text = "What we propose"
    s2 = pres.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(4))
    tb2.text_frame.text = "\n".join(f"- {b}" for b in body)

    s3 = pres.slides.add_slide(blank)
    tb3 = s3.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
    tb3.text_frame.text = cta

    out_path = OUT_DIR / f"{lead_id or 'pitch'}.pptx"
    pres.save(str(out_path))
    return out_path


@tool
def build_pptx(deck_json: dict, lead_id: str = "", brand: dict | None = None) -> dict[str, Any]:
    """Render a pitch deck (.pptx) from a deck spec.

    Args:
        deck_json: the deck spec (see codex_brain pptx_renderer for schema).
        lead_id: lead id for filename + audit row.
        brand: optional brand dict {primary, accent, name}.

    Returns:
        {ok: bool, data: {path}, error: str|None, mode: 'full'|'minimal'}
    """
    actor = get_actor_name()
    db = get_sync_db()
    lead = lead_id or "pitch"

    try:
        if _render_pptx_full is not None:
            path = _render_pptx_full(deck_json or {}, brand or {}, lead_id=lead, out_dir=OUT_DIR)
            mode = "full"
        else:
            path = _render_minimal(deck_json or {}, lead_id=lead)
            mode = "minimal"
    except Exception as e:  # noqa: BLE001
        log.warning("pptx render failed: %s", type(e).__name__)
        write_audit_sync(
            db=db,
            action="swarm.pptx.build",
            actor=actor,
            lead_id=lead_id or None,
            metadata={"error": type(e).__name__},
        )
        return {"ok": False, "data": None, "error": type(e).__name__, "mode": "none"}

    write_audit_sync(
        db=db,
        action="swarm.pptx.build",
        actor=actor,
        lead_id=lead_id or None,
        cost_cents=0,
        metadata={"path": str(path), "mode": mode},
    )
    return {"ok": True, "data": {"path": str(path)}, "error": None, "mode": mode}
