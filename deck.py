"""
SolarReach Codex Brain — Pitch deck generator
Sonnet 4.6 → JSON → python-pptx PPTX → libreoffice PDF
"""

from __future__ import annotations

import json
import subprocess
import tempfile
from pathlib import Path
from dataclasses import dataclass
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

from codex_brain.anthropic_client import get_client, SONNET, UsageRecord
from codex_brain.constants_funding import FUNDING_MODELS

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

SOLAR_ORANGE = RGBColor(0xE8, 0x5D, 0x2B)
DARK_BG      = RGBColor(0x12, 0x12, 0x12)
OFF_WHITE    = RGBColor(0xF5, 0xF3, 0xEF)
MID_GREY     = RGBColor(0x88, 0x87, 0x80)
SLIDE_W      = Inches(13.33)
SLIDE_H      = Inches(7.5)


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class LeadBrief:
    company_name: str
    decision_maker_name: str
    decision_maker_role: str
    postcode: str
    roof_area_m2: float
    annual_kwh: float
    panel_count: int
    payback_years: float
    pitch_theme: str = "GRID INDEPENDENCE"
    client_name: str = "GreenSolar UK"

    @property
    def annual_saving_gbp(self) -> float:
        return round(self.annual_kwh * 0.28)

    @property
    def capex_gbp(self) -> float:
        return round(self.panel_count * 350)

    @property
    def npv_25yr_gbp(self) -> float:
        return round(self.annual_saving_gbp * 25 * 0.85 - self.capex_gbp)


@dataclass
class PitchResult:
    slides: list[dict]
    emails: list[dict]
    funding_models: list[dict]
    pptx_path: Path
    pdf_path: Path | None
    usage: UsageRecord


# ---------------------------------------------------------------------------
# System prompt
# ---------------------------------------------------------------------------

_SYSTEM = """\
You are SolarReach Codex Brain, the AI content engine for a UK commercial solar sales platform.
Your outputs feed directly into python-pptx to generate 16:9 PPTX slide decks.

Tone: confident, data-driven, concise. No fluff, no buzzwords.
Theme: {theme}. Frame all benefits as grid independence, cost certainty, sustainability.

Output ONLY valid JSON — no markdown fences, no commentary, no preamble.
"""

_USER_TMPL = """\
Generate a complete pitch deck for:

Company:        {company_name}
Decision-maker: {dm_name}, {dm_role}
Postcode:       {postcode}
Roof area:      {area} m²
Annual yield:   {kwh} kWh
Panels:         {panels}
Annual saving:  £{saving:,}
System capex:   £{capex:,}
Payback:        {payback} years
25yr NPV:       £{npv:,}
Presenting as:  {client_name}

Return exactly this JSON shape:
{{
  "slides": [
    {{
      "title": "...",
      "subtitle": "...(optional, one line)",
      "bullets": ["...", "..."],
      "speaker_note": "...(1-2 sentences)"
    }}
  ],
  "emails": [
    {{
      "subject": "...",
      "body": "...(3-4 sentences, Variant A)"
    }},
    {{
      "subject": "...",
      "body": "...(3-4 sentences, Variant B — different angle)"
    }}
  ]
}}

Slides must be exactly 11 slides in this order:
1. Title slide
2. The energy cost problem
3. Why solar, why now
4. Your rooftop opportunity (use the real numbers above)
5. Solar radiance analysis (describe the roof data)
6. Financial model — headline ROI
7. Funding options (reference Capital Expense / Free Install / Lease / HP)
8. Our installation process (4 steps)
9. Environmental impact
10. Case study / social proof
11. Next steps & call to action
"""


# ---------------------------------------------------------------------------
# Generator
# ---------------------------------------------------------------------------

def generate_pitch(brief: LeadBrief, output_dir: Path) -> PitchResult:
    """
    Full pipeline:
      1. Call Sonnet 4.6 with cached system prompt → parse JSON
      2. Build matplotlib ROI chart
      3. Render python-pptx PPTX
      4. Convert to PDF via libreoffice headless
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    client = get_client()

    system = _SYSTEM.format(theme=brief.pitch_theme)
    user = _USER_TMPL.format(
        company_name=brief.company_name,
        dm_name=brief.decision_maker_name,
        dm_role=brief.decision_maker_role,
        postcode=brief.postcode,
        area=int(brief.roof_area_m2),
        kwh=int(brief.annual_kwh),
        panels=brief.panel_count,
        saving=int(brief.annual_saving_gbp),
        capex=int(brief.capex_gbp),
        payback=brief.payback_years,
        npv=int(brief.npv_25yr_gbp),
        client_name=brief.client_name,
    )

    raw, usage = client.complete(
        system=system,
        user=user,
        model=SONNET,
        max_tokens=3000,
        cache_system=True,
        call_type="pitch_deck",
    )

    data = json.loads(raw)
    slides: list[dict] = data["slides"]
    emails: list[dict] = data["emails"]

    # Build ROI chart PNG
    chart_path = output_dir / "roi_chart.png"
    _render_roi_chart(brief, chart_path)

    # Build PPTX
    slug = brief.company_name.lower().replace(" ", "_")[:30]
    pptx_path = output_dir / f"{slug}_pitch.pptx"
    _render_pptx(slides, brief, chart_path, pptx_path)

    # Convert to PDF
    pdf_path = _convert_to_pdf(pptx_path, output_dir)

    return PitchResult(
        slides=slides,
        emails=emails,
        funding_models=FUNDING_MODELS,
        pptx_path=pptx_path,
        pdf_path=pdf_path,
        usage=usage,
    )


# ---------------------------------------------------------------------------
# ROI chart
# ---------------------------------------------------------------------------

def _render_roi_chart(brief: LeadBrief, out_path: Path) -> None:
    years = np.arange(0, 26)
    cumulative = -brief.capex_gbp + brief.annual_saving_gbp * years

    fig, ax = plt.subplots(figsize=(8, 4), facecolor="#121212")
    ax.set_facecolor("#1A1A1A")

    # Shade negative / positive
    ax.fill_between(years, cumulative, 0, where=(cumulative < 0),
                    alpha=0.3, color="#E85D2B", label="_neg")
    ax.fill_between(years, cumulative, 0, where=(cumulative >= 0),
                    alpha=0.2, color="#3B9A1A", label="_pos")

    ax.plot(years, cumulative, color="#E85D2B", lw=2.5, zorder=3)

    # Payback marker
    pb = brief.payback_years
    ax.axvline(pb, color="#F5F3EF", lw=1, ls="--", alpha=0.6)
    ax.text(pb + 0.3, cumulative.min() * 0.85, f"Payback yr {pb:.0f}",
            color="#F5F3EF", fontsize=8, alpha=0.9)

    ax.axhline(0, color="#555", lw=0.8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x, _: f"£{int(x/1000):,}k"))
    ax.set_xlabel("Year", color="#888")
    ax.set_ylabel("Cumulative cash flow", color="#888")
    ax.tick_params(colors="#888")
    for spine in ax.spines.values():
        spine.set_edgecolor("#333")

    ax.set_title(f"25-year ROI — {brief.company_name}",
                 color="#F5F3EF", fontsize=11, pad=10)

    plt.tight_layout()
    fig.savefig(out_path, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------

def _render_pptx(
    slides: list[dict],
    brief: LeadBrief,
    chart_path: Path,
    out_path: Path,
) -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        _fill_bg(slide, DARK_BG)

        is_title_slide = i == 0
        if is_title_slide:
            _add_title_slide(slide, slide_data, brief)
        elif "financial" in slide_data.get("title", "").lower() or i == 5:
            _add_chart_slide(slide, slide_data, chart_path)
        else:
            _add_content_slide(slide, slide_data, i + 1, len(slides))

    prs.save(out_path)


def _fill_bg(slide: Any, color: RGBColor) -> None:
    from pptx.util import Emu
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(slide: Any, data: dict, brief: LeadBrief) -> None:
    # Orange accent bar on left
    _add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, SOLAR_ORANGE)

    # Title
    tf = _add_textbox(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.6))
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "Commercial Solar Opportunity")
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = OFF_WHITE

    # Subtitle
    if data.get("subtitle"):
        tf2 = _add_textbox(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
        p2 = tf2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = data["subtitle"]
        run2.font.size = Pt(20)
        run2.font.color.rgb = MID_GREY

    # Company name
    tf3 = _add_textbox(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.6))
    p3 = tf3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"Prepared for {brief.company_name} · {brief.postcode}"
    run3.font.size = Pt(13)
    run3.font.color.rgb = SOLAR_ORANGE

    # Presenting-as
    tf4 = _add_textbox(slide, Inches(0.5), Inches(6.8), Inches(6), Inches(0.4))
    p4 = tf4.text_frame.paragraphs[0]
    run4 = p4.add_run()
    run4.text = brief.client_name
    run4.font.size = Pt(11)
    run4.font.color.rgb = MID_GREY


def _add_content_slide(slide: Any, data: dict, slide_num: int, total: int) -> None:
    # Thin orange top bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), SOLAR_ORANGE)

    # Slide number bottom-right
    tf_num = _add_textbox(slide, Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.3))
    p_num = tf_num.text_frame.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    run_num = p_num.add_run()
    run_num.text = f"{slide_num} / {total}"
    run_num.font.size = Pt(9)
    run_num.font.color.rgb = MID_GREY

    # Title
    tf_title = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.9))
    p_title = tf_title.text_frame.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = data.get("title", "")
    run_title.font.bold = True
    run_title.font.size = Pt(28)
    run_title.font.color.rgb = OFF_WHITE

    # Subtitle
    y_off = Inches(1.25)
    if data.get("subtitle"):
        tf_sub = _add_textbox(slide, Inches(0.6), y_off, Inches(11), Inches(0.5))
        p_sub = tf_sub.text_frame.paragraphs[0]
        run_sub = p_sub.add_run()
        run_sub.text = data["subtitle"]
        run_sub.font.size = Pt(15)
        run_sub.font.color.rgb = SOLAR_ORANGE
        y_off = Inches(1.8)

    # Bullets
    bullets: list[str] = data.get("bullets", [])
    if bullets:
        tf_b = _add_textbox(slide, Inches(0.8), y_off, Inches(11), Inches(7.5) - y_off - Inches(0.6))
        tf_b.text_frame.word_wrap = True
        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf_b.text_frame.paragraphs[0]
            else:
                p = tf_b.text_frame.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(16)
            run.font.color.rgb = OFF_WHITE


def _add_chart_slide(slide: Any, data: dict, chart_path: Path) -> None:
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), SOLAR_ORANGE)

    # Title
    tf = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.9))
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "Financial model")
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = OFF_WHITE

    # Insert chart image
    if chart_path.exists():
        slide.shapes.add_picture(
            str(chart_path),
            Inches(0.5), Inches(1.4),
            Inches(12.3), Inches(5.5),
        )


def _add_rect(slide: Any, left: Emu, top: Emu, width: Emu, height: Emu, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_textbox(slide: Any, left: Emu, top: Emu, width: Emu, height: Emu) -> Any:
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox


# ---------------------------------------------------------------------------
# PDF conversion
# ---------------------------------------------------------------------------

def _convert_to_pdf(pptx_path: Path, out_dir: Path) -> Path | None:
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                str(pptx_path),
                "--outdir", str(out_dir),
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode != 0:
            return None
        pdf_path = out_dir / (pptx_path.stem + ".pdf")
        return pdf_path if pdf_path.exists() else None
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return None
